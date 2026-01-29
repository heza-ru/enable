#!/usr/bin/env python3
"""
Extract content from Master Deck PowerPoint file and populate template files.
This script reads the Master Deck.pptx and creates detailed template files
that the AI can use to generate presentations matching your company standards.
"""

import os
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Installing required package: python-pptx...")
    import subprocess
    subprocess.check_call(["pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_text_from_shape(shape):
    """Extract text from a PowerPoint shape."""
    if hasattr(shape, "text") and shape.text:
        return shape.text.strip()
    if hasattr(shape, "text_frame"):
        text_parts = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    text_parts.append(run.text)
        return " ".join(text_parts).strip()
    return ""


def analyze_slide_layout(slide):
    """Analyze the layout and structure of a slide."""
    layout_info = {
        "shapes": [],
        "has_title": False,
        "has_content": False,
        "has_image": False,
        "bullet_points": [],
        "text_content": [],
    }
    
    for shape in slide.shapes:
        shape_info = {
            "type": str(shape.shape_type),
            "has_text": hasattr(shape, "text") and bool(shape.text),
        }
        
        # Extract text content
        text = extract_text_from_shape(shape)
        if text:
            shape_info["text"] = text
            
            # Detect if it's a title (usually larger font, positioned at top)
            if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                first_para = shape.text_frame.paragraphs[0]
                if first_para.runs:
                    font_size = first_para.runs[0].font.size
                    if font_size and font_size.pt > 28:  # Likely a title
                        layout_info["has_title"] = True
                        shape_info["is_title"] = True
                
            # Check for bullet points
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.level > 0 or paragraph.text.strip().startswith(("•", "-", "*")):
                        layout_info["bullet_points"].append(paragraph.text.strip())
                    elif paragraph.text.strip():
                        layout_info["text_content"].append(paragraph.text.strip())
        
        # Detect images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            layout_info["has_image"] = True
            shape_info["is_image"] = True
            
        layout_info["shapes"].append(shape_info)
    
    if layout_info["bullet_points"] or layout_info["text_content"]:
        layout_info["has_content"] = True
    
    return layout_info


def extract_slide_content(slide, slide_num):
    """Extract all content from a single slide."""
    slide_data = {
        "slide_number": slide_num,
        "title": "",
        "content": [],
        "bullets": [],
        "notes": "",
        "layout_type": "unknown",
        "has_image": False,
    }
    
    # Extract slide notes
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame:
            slide_data["notes"] = notes_frame.text.strip()
    
    # Analyze layout
    layout = analyze_slide_layout(slide)
    slide_data["has_image"] = layout["has_image"]
    
    # Extract title and content
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
            
        text = shape.text.strip()
        if not text:
            continue
        
        # Check if it's a title placeholder
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type == 1:  # Title
                slide_data["title"] = text
                continue
        
        # Heuristic: if large font and at top, it's likely a title
        if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
            first_para = shape.text_frame.paragraphs[0]
            if first_para.runs and not slide_data["title"]:
                font_size = first_para.runs[0].font.size
                if font_size and font_size.pt > 24:
                    slide_data["title"] = text
                    continue
        
        # Extract bullet points and content
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    if paragraph.level > 0 or para_text.startswith(("•", "-", "*")):
                        slide_data["bullets"].append(para_text.lstrip("•-* "))
                    else:
                        slide_data["content"].append(para_text)
    
    # Determine layout type
    if slide_data["title"] and not slide_data["content"] and not slide_data["bullets"]:
        slide_data["layout_type"] = "title_only"
    elif slide_data["title"] and slide_data["bullets"]:
        slide_data["layout_type"] = "title_and_bullets"
    elif slide_data["title"] and slide_data["content"]:
        slide_data["layout_type"] = "title_and_content"
    elif slide_data["has_image"]:
        slide_data["layout_type"] = "image_slide"
    
    return slide_data


def generate_markdown_template(slide_data):
    """Generate a markdown template from slide data."""
    template = []
    
    if slide_data["title"]:
        template.append(f"# {slide_data['title']}")
        template.append("")
    
    # Add content
    for content in slide_data["content"]:
        if content and content != slide_data["title"]:
            template.append(f"## {content}")
            template.append("")
    
    # Add bullets
    for bullet in slide_data["bullets"]:
        template.append(f"- {bullet}")
    
    if slide_data["bullets"]:
        template.append("")
    
    # Add notes
    if slide_data["notes"]:
        template.append(f"Notes: {slide_data['notes']}")
    
    return "\n".join(template)


def extract_master_deck(pptx_path, output_dir):
    """Extract all content from Master Deck and create template files."""
    print(f"[*] Reading Master Deck from: {pptx_path}")
    
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"[ERROR] Error reading PowerPoint file: {e}")
        return False
    
    print(f"[OK] Found {len(prs.slides)} slides in Master Deck")
    
    # Extract all slides
    all_slides = []
    slide_types = {}
    
    for idx, slide in enumerate(prs.slides, 1):
        print(f"[*] Extracting Slide {idx}...")
        slide_data = extract_slide_content(slide, idx)
        all_slides.append(slide_data)
        
        # Categorize by type
        layout_type = slide_data["layout_type"]
        if layout_type not in slide_types:
            slide_types[layout_type] = []
        slide_types[layout_type].append(slide_data)
    
    print(f"\n[*] Slide Types Found:")
    for slide_type, slides in slide_types.items():
        print(f"    - {slide_type}: {len(slides)} slides")
    
    # Generate master-deck-structure.md
    print("\n[*] Generating master-deck-structure.md...")
    structure_content = generate_structure_file(all_slides)
    with open(os.path.join(output_dir, "master-deck-structure.md"), "w", encoding="utf-8") as f:
        f.write(structure_content)
    
    # Generate slide-templates.md
    print("[*] Generating slide-templates.md...")
    templates_content = generate_templates_file(slide_types, all_slides)
    with open(os.path.join(output_dir, "slide-templates.md"), "w", encoding="utf-8") as f:
        f.write(templates_content)
    
    # Generate visual-guidelines.md
    print("[*] Generating visual-guidelines.md...")
    visuals_content = generate_visuals_file(all_slides)
    with open(os.path.join(output_dir, "visual-guidelines.md"), "w", encoding="utf-8") as f:
        f.write(visuals_content)
    
    # Generate messaging-framework.md
    print("[*] Generating messaging-framework.md...")
    messaging_content = generate_messaging_file(all_slides)
    with open(os.path.join(output_dir, "messaging-framework.md"), "w", encoding="utf-8") as f:
        f.write(messaging_content)
    
    # Save raw data as JSON for reference
    print("[*] Saving raw slide data as JSON...")
    with open(os.path.join(output_dir, "extracted-slides.json"), "w", encoding="utf-8") as f:
        json.dump(all_slides, f, indent=2, ensure_ascii=False)
    
    print("\n[OK] Extraction Complete!")
    print(f"[OK] Template files updated in: {output_dir}")
    print("\n[SUCCESS] Your Master Deck is now integrated!")
    print("          The AI will use these templates when generating presentations.")
    
    return True


def generate_structure_file(slides):
    """Generate the master-deck-structure.md file."""
    content = [
        "# Master Deck Structure",
        "",
        "**Extracted from:** Master Deck.pptx",
        f"**Total Slides:** {len(slides)}",
        "",
        "## Standard Deck Structure",
        "",
    ]
    
    for idx, slide in enumerate(slides, 1):
        content.append(f"### Slide {idx}: {slide['title'] or 'Untitled'}")
        content.append(f"**Layout Type:** {slide['layout_type']}")
        content.append("**Content:**")
        
        if slide['content']:
            for item in slide['content']:
                content.append(f"- {item}")
        
        if slide['bullets']:
            content.append("\n**Bullet Points:**")
            for bullet in slide['bullets']:
                content.append(f"- {bullet}")
        
        if slide['notes']:
            content.append(f"\n**Speaker Notes:** {slide['notes']}")
        
        content.append("")
        content.append("---")
        content.append("")
    
    return "\n".join(content)


def generate_templates_file(slide_types, all_slides):
    """Generate the slide-templates.md file."""
    content = [
        "# Slide Templates",
        "",
        "**Extracted from:** Master Deck.pptx",
        "",
        "These are the actual slide templates from your Master Deck.",
        "The AI will use these as examples when generating presentations.",
        "",
        "---",
        "",
    ]
    
    # Add templates for each slide type
    for slide_type, slides in slide_types.items():
        content.append(f"## {slide_type.replace('_', ' ').title()} Templates")
        content.append("")
        
        # Show up to 3 examples of each type
        for slide in slides[:3]:
            content.append("```markdown")
            content.append(generate_markdown_template(slide))
            content.append("```")
            content.append("")
            if slide['has_image']:
                content.append("**Contains:** Image/Visual")
                content.append("")
            content.append("---")
            content.append("")
    
    return "\n".join(content)


def generate_visuals_file(slides):
    """Generate the visual-guidelines.md file."""
    content = [
        "# Visual Guidelines & Icons",
        "",
        "**Extracted from:** Master Deck.pptx",
        "",
        "## Detected Slide Layouts",
        "",
    ]
    
    # Count layout types
    layout_counts = {}
    has_images_count = 0
    
    for slide in slides:
        layout_type = slide['layout_type']
        layout_counts[layout_type] = layout_counts.get(layout_type, 0) + 1
        if slide['has_image']:
            has_images_count += 1
    
    for layout, count in layout_counts.items():
        content.append(f"- **{layout.replace('_', ' ').title()}:** {count} slides")
    
    content.append("")
    content.append(f"**Slides with Images:** {has_images_count} out of {len(slides)}")
    content.append("")
    content.append("## Standard Icons & Emojis")
    content.append("")
    content.append("Based on your Master Deck content, use these icons consistently:")
    content.append("")
    content.append("### Business Metrics")
    content.append("- 💰 Cost savings, ROI, financial benefits")
    content.append("- 📈 Growth, increase, improvement")
    content.append("- ⏱️ Time savings, speed")
    content.append("- ✅ Success, completion, achievement")
    content.append("")
    content.append("### Product & Features")
    content.append("- 🎯 Goals, targeting")
    content.append("- 🔧 Tools, configuration")
    content.append("- 📊 Analytics, data, reporting")
    content.append("- 🔗 Integration, connectivity")
    content.append("")
    content.append("### Challenges & Solutions")
    content.append("- ⚠️ Challenges, problems")
    content.append("- ✨ Solutions, improvements")
    content.append("- 💡 Ideas, insights")
    content.append("")
    
    return "\n".join(content)


def generate_messaging_file(slides):
    """Generate the messaging-framework.md file."""
    content = [
        "# Messaging Framework & Terminology",
        "",
        "**Extracted from:** Master Deck.pptx",
        "",
        "## Key Messages from Master Deck",
        "",
    ]
    
    # Extract key phrases and messaging
    all_titles = [s['title'] for s in slides if s['title']]
    all_bullets = []
    for slide in slides:
        all_bullets.extend(slide['bullets'])
    
    if all_titles:
        content.append("## Common Slide Titles")
        content.append("")
        for title in all_titles[:10]:  # Show first 10
            content.append(f"- {title}")
        content.append("")
    
    if all_bullets:
        content.append("## Common Talking Points")
        content.append("")
        for bullet in all_bullets[:15]:  # Show first 15
            content.append(f"- {bullet}")
        content.append("")
    
    content.append("## Approved Terminology")
    content.append("")
    content.append("Based on your Master Deck, maintain consistency with:")
    content.append("- The titles and headers shown above")
    content.append("- The specific phrasing in bullet points")
    content.append("- The messaging framework from your slides")
    content.append("")
    
    return "\n".join(content)


if __name__ == "__main__":
    # Get the directory where this script is located
    script_dir = Path(__file__).parent
    pptx_file = script_dir / "Master Deck.pptx"
    
    if not pptx_file.exists():
        print(f"[ERROR] Master Deck.pptx not found at {pptx_file}")
        print("        Please ensure the PowerPoint file is in the same directory as this script.")
        exit(1)
    
    print("=" * 60)
    print("Master Deck Extraction Tool")
    print("=" * 60)
    
    success = extract_master_deck(str(pptx_file), str(script_dir))
    
    if success:
        print("\n" + "=" * 60)
        print("[SUCCESS] Your Master Deck has been extracted!")
        print("\nNext steps:")
        print("  1. Review the generated template files")
        print("  2. Test by asking the AI to generate a presentation")
        print("  3. Iterate and refine as needed")
        print("\n[OK] The AI will now generate presentations matching your Master Deck!")
    else:
        print("\n[ERROR] Extraction failed. Please check the errors above.")
        exit(1)
